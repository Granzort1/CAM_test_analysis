from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import cv2

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)

    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

def get_video_frame(video_path):
    video = cv2.VideoCapture(video_path)
    success, image = video.read()
    if success:
        return cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
    return None

def add_video_slide(prs, title, air_video_path, soil_video_path):
    if not (os.path.exists(air_video_path) and os.path.exists(soil_video_path)):
        print(f"Skipping slide for {title} due to missing video files")
        return

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 제목 설정
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)

    # 슬라이드 크기
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 비디오 삽입 함수
    def insert_video(video_path, left, top, width, height):
        poster_frame = get_video_frame(video_path)
        if poster_frame is not None:
            temp_image_path = 'temp_poster_frame.png'
            cv2.imwrite(temp_image_path, cv2.cvtColor(poster_frame, cv2.COLOR_RGB2BGR))
            slide.shapes.add_movie(video_path, left, top, width, height,
                                   mime_type='video/mp4', poster_frame_image=temp_image_path)
            os.remove(temp_image_path)
        else:
            slide.shapes.add_movie(video_path, left, top, width, height, mime_type='video/mp4')

    # Air 비디오 추가
    video_width = (slide_width - Inches(0.75)) / 2
    video_height = slide_height - title_shape.height - Inches(1.25)
    insert_video(air_video_path, Inches(0.25), title_shape.height + Inches(0.25), video_width, video_height)

    # Soil 비디오 추가
    insert_video(soil_video_path, slide_width / 2 + Inches(0.125), title_shape.height + Inches(0.25), video_width, video_height)

    # 캡션 추가
    left = Inches(0.25)
    top = slide_height - Inches(0.7)
    width = slide_width - Inches(0.5)
    height = Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"Left: Air Concentration Animation, Right: Soil Concentration Animation"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

def create_presentation_with_videos():
    # PowerPoint 파일 경로
    ppt_file = r'C:\CAM_test_analysis\graph\Concentration_Animations_Analysis.pptx'

    # 동영상이 저장된 폴더 경로
    video_folder = r'C:\CAM_test_analysis\animations'

    # 물질 목록
    substances = [
        "Ethylacetate", "Benzene", "Methylacrylate", "Methyltrichlorosilane", "Ethyleneoxide",
        "Triethylamine", "Methylethylketoneperoxide", "Methylhydrazine", "Chloromethane", "Methylamine",
        "Vinylchloride", "Carbondisulfide", "Trimethylamine", "Propyleneoxide", "Methylvinylketone", "Nitrobenzene"
    ]

    # 새 프레젠테이션 생성
    prs = Presentation()

    # 슬라이드 크기 설정 (16:9 비율)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 제목 슬라이드 추가
    add_title_slide(prs,
                    "Concentration Analysis of Various Substances",
                    "Air and Soil Concentration Animations")

    # 각 물질에 대한 슬라이드 생성
    for i in range(26, 42):
        substance_index = i - 26
        if substance_index < len(substances):
            substance_name = substances[substance_index]
        else:
            substance_name = f"Unknown Substance {i}"

        air_video_path = os.path.join(video_folder, f'Concentration{i}_Air_animation.mp4')
        soil_video_path = os.path.join(video_folder, f'Concentration{i}_Soil_animation.mp4')

        add_video_slide(prs,
                        f"{i}. Concentration Animations for {substance_name}",
                        air_video_path,
                        soil_video_path)

    # PowerPoint 파일 저장
    prs.save(ppt_file)

    print(f"PowerPoint 프레젠테이션이 '{ppt_file}'에 저장되었습니다.")

if __name__ == "__main__":
    create_presentation_with_videos()