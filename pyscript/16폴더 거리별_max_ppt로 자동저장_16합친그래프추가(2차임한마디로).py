from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from PIL import Image

# PowerPoint 파일 경로
ppt_file = r'C:\CAM_test_analysis\output\Concentration_Graphs_Analysis.pptx'

# 그래프가 저장된 폴더 경로
graph_folder = r'C:\CAM_test_analysis\output'

# 물질 목록
substances = [
    "Ethylacetate", "Benzene", "Methylacrylate", "Methyltrichlorosilane", "Ethyleneoxide",
    "Triethylamine", "Methylethylketoneperoxide", "Methylhydrazine", "Chloromethane", "Methylamine",
    "Vinylchloride", "Carbondisulfide", "Trimethylamine", "Propyleneoxide", "Methylvinylketone", "Nitrobenzene"
]

# 거리 구간 정의 (미터 단위)
distance_ranges = [(0, 500), (500, 1000), (1000, 3000), (3000, 5000), (5000, 7000)]


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)

    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)


def add_content_slide(prs, title, air_img_path, soil_img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 제목 설정
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)

    # 슬라이드 크기
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 이미지 삽입 함수
    def insert_image(img_path, left):
        if os.path.exists(img_path):
            with Image.open(img_path) as img:
                img_width, img_height = img.size

            # 이미지 크기 계산 (원본 비율 유지)
            max_height = slide_height - title_shape.height - Inches(0.5)
            max_width = (slide_width - Inches(0.75)) / 2

            width = max_width
            height = (max_width / img_width) * img_height

            if height > max_height:
                height = max_height
                width = (max_height / img_height) * img_width

            top = title_shape.height + Inches(0.25)
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)

    # Air 그래프 추가
    insert_image(air_img_path, Inches(0.25))

    # Soil 그래프 추가
    insert_image(soil_img_path, slide_width / 2 + Inches(0.125))

    # 캡션 추가
    left = Inches(0.25)
    top = slide_height - Inches(0.7)
    width = slide_width - Inches(0.5)
    height = Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"Left: Air Concentration, Right: Soil Concentration"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER


# 새 프레젠테이션 생성
prs = Presentation()

# 슬라이드 크기 설정 (16:9 비율)
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# 제목 슬라이드 추가
add_title_slide(prs,
                "Concentration Analysis of Various Substances",
                "Air and Soil Concentration Graphs")

# 각 물질에 대한 슬라이드 생성
for i in range(26, 42):
    substance_index = i - 26
    if substance_index < len(substances):
        substance_name = substances[substance_index]
    else:
        substance_name = f"Unknown Substance {i}"

    air_img_path = os.path.join(graph_folder, f'{i}_Air.png')
    soil_img_path = os.path.join(graph_folder, f'{i}_Soil.png')

    add_content_slide(prs,
                      f"{i}. Concentration Graphs for {substance_name}",
                      air_img_path,
                      soil_img_path)

# 거리별 대기와 토양 그래프 슬라이드 추가
for start, end in distance_ranges:
    air_img_path = os.path.join(graph_folder, f'Air_{start}m-{end}m.png')
    soil_img_path = os.path.join(graph_folder, f'Soil_{start}m-{end}m.png')

    add_content_slide(prs,
                      f"Concentration Graphs for {start}m-{end}m Range",
                      air_img_path,
                      soil_img_path)

# PowerPoint 파일 저장
prs.save(ppt_file)

print(f"PowerPoint 프레젠테이션이 '{ppt_file}'에 저장되었습니다.")